import os
import sys
import glob
import re
from pptx import Presentation
from pptx.util import Inches
from google.adk.tools.tool_context import ToolContext
from google.genai.types import Part

try:
    from ..config import save_artifact_helper, get_gcs_artifact_url
    from ..tools.file_manager import get_topic_slug
except ImportError:
    from config import save_artifact_helper, get_gcs_artifact_url
    from tools.file_manager import get_topic_slug


def extract_script(md_path: str) -> str:
    if not os.path.exists(md_path):
        return ""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    match = re.search(r'## Script\s*([\s\S]*)', content)
    if match:
        return match.group(1).strip()
    return ""


def export_session_to_pptx(session_path: str, pptx_file_name: str = None) -> dict:
    """Gathers all slide PNGs and compiles them into a single widescreen PPTX presentation with speaker notes."""
    slides_dir = os.path.join(session_path, 'slides')
    png_files = sorted(glob.glob(os.path.join(slides_dir, 'slide_*.png')))
    
    if not png_files:
        slides_dir = session_path
        png_files = sorted(glob.glob(os.path.join(slides_dir, 'slide_*.png')))
        
    if not png_files:
        raise ValueError(f"No slide PNG images found in {session_path}. Make sure to generate slide images first.")
        
    output_name = pptx_file_name or 'presentation.pptx'
    if not output_name.endswith('.pptx'):
        output_name += '.pptx'
        
    output_path = os.path.join(session_path, output_name)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    for png_file in png_files:
        filename = os.path.basename(png_file)
        slide_num_match = re.search(r'slide_(\d+)\.png', filename)
        if not slide_num_match:
            continue
        pad_num = slide_num_match.group(1)
        
        md_file = os.path.join(session_path, f"slide_{pad_num}.md")
        script_notes = extract_script(md_file)
        
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(png_file, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        
        if script_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = script_notes
            
    prs.save(output_path)
    
    return {
        "message": f"Successfully compiled {len(png_files)} slides into a widescreen PPTX with speaker notes.",
        "pptxName": output_name,
        "pptxPath": output_path
    }


async def export_deck_pptx(session_path: str, tool_context: ToolContext) -> str:
    """Compiles all generated slide PNG images in the active session into a single PPTX presentation,
    attaching the corresponding speaker notes to each slide.
    
    Args:
        session_path: The absolute session path returned by initialize_session
        tool_context: The tool context injected by the framework
    """
    slug = get_topic_slug(session_path)
    pptx_filename = f"{slug}.pptx"
    
    try:
        result = export_session_to_pptx(session_path, pptx_filename)
        pptx_path = result["pptxPath"]
        
        with open(pptx_path, 'rb') as f:
            pptx_bytes = f.read()
            
        artifact_part = Part.from_bytes(
            data=pptx_bytes,
            mime_type="application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )
        version = await save_artifact_helper(pptx_filename, artifact_part, tool_context)

        gcs_url = get_gcs_artifact_url(pptx_filename, tool_context, version=version)
        if gcs_url:
            return f"Presentation PPTX successfully compiled with speaker notes.\nDownload it here: {gcs_url}"

        return f"Presentation PPTX successfully compiled and saved to {pptx_path}"
    except Exception as e:
        return f"Failed to export PPTX: {str(e)}"
