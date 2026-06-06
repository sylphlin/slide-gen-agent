import os
import sys
import glob
import re
import json
from pptx import Presentation
from pptx.util import Inches

def extract_script(md_path: str) -> str:
    if not os.path.exists(md_path):
        return ""
    with open(md_path, 'r', encoding='utf-8') as f:
        content = f.read()
    match = re.search(r'## Script\s*([\s\S]*)', content)
    if match:
        return match.group(1).strip()
    return ""

def export_session_to_pptx(session_path: str, pptx_file_name: str = None) -> dict:
    """Gathers all slide PNGs and compiles them into a single widescreen PPTX presentation with speaker notes."""
    slides_dir = os.path.join(session_path, 'slides')
    png_files = sorted(glob.glob(os.path.join(slides_dir, 'slide_*.png')))
    
    if not png_files:
        slides_dir = session_path
        png_files = sorted(glob.glob(os.path.join(slides_dir, 'slide_*.png')))
        
    if not png_files:
        raise ValueError(f"No slide PNG images found in {session_path}. Make sure to generate slide images first.")
        
    output_name = pptx_file_name or 'presentation.pptx'
    if not output_name.endswith('.pptx'):
        output_name += '.pptx'
        
    output_path = os.path.join(session_path, output_name)
    
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]
    
    for png_file in png_files:
        filename = os.path.basename(png_file)
        slide_num_match = re.search(r'slide_(\d+)\.png', filename)
        if not slide_num_match:
            continue
        pad_num = slide_num_match.group(1)
        
        md_file = os.path.join(session_path, f"slide_{pad_num}.md")
        script_notes = extract_script(md_file)
        
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(png_file, Inches(0), Inches(0), width=prs.slide_width, height=prs.slide_height)
        
        if script_notes:
            notes_slide = slide.notes_slide
            text_frame = notes_slide.notes_text_frame
            text_frame.text = script_notes
            
    prs.save(output_path)
    
    return {
        "message": f"Successfully compiled {len(png_files)} slides into a widescreen PPTX with speaker notes.",
        "pptxName": output_name,
        "pptxPath": output_path
    }

if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python3 pptx_exporter.py <sessionPath> [pptxFileName]", file=sys.stderr)
        sys.exit(1)
        
    session_path_arg = sys.argv[1]
    pptx_name_arg = sys.argv[2] if len(sys.argv) > 2 else None
    
    try:
        result = export_session_to_pptx(session_path_arg, pptx_name_arg)
        print(json.dumps(result, indent=2))
        sys.exit(0)
    except Exception as e:
        print(f"PPTX compilation failed: {str(e)}", file=sys.stderr)
        sys.exit(1)
